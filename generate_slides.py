"""Generate a sample progress presentation for CivilQntify."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color palette
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)    # Blue
SECONDARY = RGBColor(0x05, 0x96, 0x69)  # Green
ACCENT = RGBColor(0x7C, 0x3A, 0xED)     # Purple
DARK = RGBColor(0x1F, 0x29, 0x37)       # Dark gray
LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)   # Light gray background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, color=PRIMARY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title_text(slide, text, left=Inches(0.8), top=Inches(0.4), width=Inches(11), height=Inches(0.8), color=DARK, size=32):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox


def add_subtitle(slide, text, left=Inches(0.8), top=Inches(1.2), width=Inches(10), height=Inches(0.5), color=PRIMARY, size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return txBox


def add_bullet_list(slide, items, left=Inches(0.8), top=Inches(2.0), width=Inches(11), height=Inches(4.5), size=14, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_placeholder_box(slide, left, top, width, height, label="[PLACE IMAGE HERE]"):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
    shape.line.dash_style = 2  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(40)
    return shape


def add_content_card(slide, left, top, width, height, title, items, color=PRIMARY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    p.space_after = Pt(8)

    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)


# ── Slide 1: Title ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, PRIMARY)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CivilQntify"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Automated Concrete Mix Design & Material Quantification"
p2.font.size = Pt(22)
p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "[Your Name]  |  [Course/Program]  |  [Supervisor Name]"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
p3.alignment = PP_ALIGN.CENTER

p4 = tf3.add_paragraph()
p4.text = "[Date]"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(8)


# ── Slide 2: Problem Statement ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)
add_title_text(slide, "Problem Statement")
add_subtitle(slide, "Why CivilQntify?")

items = [
    "Manual concrete mix design is time-consuming and error-prone",
    "Engineers must juggle multiple code tables (ACI 211.1, IS 10262)",
    "No integrated tool for mix design → quantification → cost estimation",
    "Material wastage and cost overruns due to inaccurate quantity takeoffs",
    "Lack of standardized, auditable calculation trails for mix approval",
]
add_bullet_list(slide, items, top=Inches(2.0))

add_placeholder_box(slide, Inches(7.5), Inches(4.5), Inches(5), Inches(2.2), "[Flowchart: Current Manual Workflow]")


# ── Slide 3: Objectives ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide, SECONDARY)
add_title_text(slide, "Project Objectives")
add_subtitle(slide, "What we aim to achieve", color=SECONDARY)

items = [
    "Automate concrete mix proportioning per ACI 211.1 and IS 10262:2019",
    "Provide material quantification (volume-based & element-based)",
    "Integrate cost estimation with local pricing (GH Cedis)",
    "Generate professional PDF/CSV/JSON reports for mix documentation",
    "Support SCM replacements, admixtures, and exposure classes",
    "Deliver a clean, tabbed desktop GUI for civil engineers",
]
add_bullet_list(slide, items, top=Inches(2.0))

add_content_card(slide, Inches(8), Inches(2.0), Inches(4.5), Inches(2.5),
    "Scope", [
        "Desktop application (PyQt6)",
        "Two design codes (ACI + IS)",
        "Offline — no server required",
    ], color=SECONDARY)


# ── Slide 4: Architecture ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide, ACCENT)
add_title_text(slide, "System Architecture")
add_subtitle(slide, "High-level component overview", color=ACCENT)

add_placeholder_box(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(5.0),
    "[INSERT: architecture_diagram.png — use the generated diagram]")


# ── Slide 5: Technical Stack ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)
add_title_text(slide, "Technical Stack")

add_content_card(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(4.5),
    "Frontend (GUI)", [
        "PyQt6 desktop framework",
        "QTabWidget navigation",
        "QSS custom design system",
        "Inter + JetBrains Mono fonts",
        "QThread background workers",
    ], color=PRIMARY)

add_content_card(slide, Inches(4.9), Inches(2.0), Inches(3.5), Inches(4.5),
    "Backend (Calculation)", [
        "Pure Python calculation engine",
        "Frozen dataclasses for immutability",
        "Strategy pattern for design codes",
        "Absolute volume method",
        "Qt signals/slots for data flow",
    ], color=SECONDARY)

add_content_card(slide, Inches(9.0), Inches(2.0), Inches(3.5), Inches(4.5),
    "Export & Reporting", [
        "PDF generation (fpdf2)",
        "CSV export",
        "JSON export",
        "HTML report preview",
        "Text report generation",
    ], color=ACCENT)


# ── Slide 6: Concrete Mix Design UI ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)
add_title_text(slide, "Concrete Mix Design — UI")
add_subtitle(slide, "3-step input form with code-aware validation")

add_placeholder_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8),
    "[INSERT: Screenshot of Concrete Mix Design tab — step 1]")

add_placeholder_box(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.8),
    "[INSERT: Screenshot of Concrete Mix Design tab — step 2 or results]")


# ── Slide 7: Calculation Engine ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide, SECONDARY)
add_title_text(slide, "Calculation Engine")
add_subtitle(slide, "Core proportioning logic", color=SECONDARY)

items = [
    "Target mean strength: f'ck = fck + 1.65 × S (IS) or f'cr = f'c + 1.34s (ACI)",
    "Water-cement ratio from strength tables with interpolation",
    "Water content selection based on slump, NMSA, and aggregate shape",
    "Absolute volume method for aggregate proportions",
    "Moisture correction for field conditions (absorption + free moisture)",
    "Cementitious efficiency factor for SCMs",
    "Exposure class limits (IS 456 / ACI 318)",
]
add_bullet_list(slide, items, top=Inches(2.0), size=13)

add_placeholder_box(slide, Inches(8), Inches(4.5), Inches(4.5), Inches(2.5),
    "[INSERT: Calculation steps tree view from result panel]")


# ── Slide 8: Design Standards ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide, ACCENT)
add_title_text(slide, "Design Standards Implemented")
add_subtitle(slide, "Strategy pattern for multi-code support", color=ACCENT)

add_content_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
    "ACI 211.1-22 (American)", [
        "Table 5.3.3 — Water content",
        "Table 6.3.3 — w/c ratio vs strength",
        "Table 5.3.6 — Coarse aggregate volume",
        "Overdesign for no production data",
        "Sulfate exposure limits (S0–S3)",
        "Air-entrained concrete support",
    ], color=PRIMARY)

add_content_card(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5),
    "IS 10262:2019 (Indian)", [
        "Table 7 — Water content",
        "Table 8 — w/c ratio vs strength",
        "Table 10 — CA fraction per zone",
        "Grading zones (I–IV)",
        "Exposure classes ( Mild to Severe )",
        "SCM replacement limits",
    ], color=SECONDARY)


# ── Slide 9: Material Quantification ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)
add_title_text(slide, "Material Quantification")
add_subtitle(slide, "From mix proportions to material bills")

items = [
    "Transfer mix design data via MixDesignTransferData (frozen dataclass)",
    "Volume-based: enter total concrete volume + wastage percentage",
    "Element-based: define structural elements (footing, column, beam, slab, wall)",
    "Automatic cement bag calculation (bag weight configurable)",
    "Material bill with cement bags, aggregate tonnes, water litres",
    "Override any material field without affecting other values",
]
add_bullet_list(slide, items, top=Inches(2.0))

add_placeholder_box(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5),
    "[INSERT: Screenshot of Material Quantification tab]")


# ── Slide 10: Cost Estimation ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)
add_title_text(slide, "Cost Estimation")
add_subtitle(slide, "Project costing with GH Cedis")

items = [
    "Material prices input (per bag, per tonne, per litre)",
    "Additional costs: labour, transport, overhead, profit, contingency",
    "Automatic cost breakdown from material bill",
    "Subtotal, contingency, and grand total calculation",
    "Project info: name, client, date, engineer, notes",
    "Tab now visible in the frontend (recently re-enabled)",
]
add_bullet_list(slide, items, top=Inches(2.0))

add_placeholder_box(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5),
    "[INSERT: Screenshot of Cost Estimation tab with results]")


# ── Slide 11: Export & Reports ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide, SECONDARY)
add_title_text(slide, "Export & Reporting")
add_subtitle(slide, "Professional documentation output", color=SECONDARY)

add_content_card(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(4.5),
    "PDF Report", [
        "Full mix design report",
        "Calculation steps included",
        "Material properties table",
        "Professional formatting",
        "fpdf2 library",
    ], color=PRIMARY)

add_content_card(slide, Inches(4.9), Inches(2.0), Inches(3.5), Inches(4.5),
    "Data Export", [
        "CSV for spreadsheet analysis",
        "JSON for API integration",
        "Text report for quick review",
        "All export formats available",
        "One-click download",
    ], color=SECONDARY)

add_content_card(slide, Inches(9.0), Inches(2.0), Inches(3.5), Inches(4.5),
    "Report Preview", [
        "HTML preview before export",
        "QTextBrowser (no WebEngine)",
        "Print-ready formatting",
        "Date/time stamped",
        "Project metadata included",
    ], color=ACCENT)


# ── Slide 12: Testing ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)
add_title_text(slide, "Testing & Validation")
add_subtitle(slide, "149 tests passing — ACI & IS pipelines verified")

add_content_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.2),
    "Test Coverage", [
        "149 unit tests across 3 test files",
        "ACI 211.1: water content, w/c ratio, air content, sulfate limits",
        "IS 10262: target strength, water content, CA fraction, grading",
        "Material quantification: volume, elements, overrides, reports",
    ], color=PRIMARY)

add_content_card(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.2),
    "Integration Tests", [
        "Full ACI pipeline: input → result → export",
        "Full IS pipeline: input → result → export",
        "End-to-end quantification flow",
        "Error handling and validation edge cases",
    ], color=SECONDARY)

add_placeholder_box(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4.7),
    "[INSERT: Terminal screenshot showing pytest output — 149 passed]")


# ── Slide 13: Challenges & Lessons ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide, ACCENT)
add_title_text(slide, "Challenges & Lessons Learned")
add_subtitle(slide, "Technical decisions and tradeoffs", color=ACCENT)

items = [
    "Frozen dataclasses → immutability prevents silent calculation corruption",
    "Strategy pattern → clean separation between ACI and IS code logic",
    "QThread workers → non-blocking UI during heavy calculations",
    "Qt signals/slots → type-safe inter-tab data transfer without tight coupling",
    "Absolute volume method → ensuring unit consistency across codes",
    "Moisture correction → bridging lab conditions to field conditions",
    "PDF generation → fpdf2 chosen over ReportLab for simplicity",
]
add_bullet_list(slide, items, top=Inches(2.0), size=13)


# ── Slide 14: Next Steps ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)
add_title_text(slide, "Next Steps & Roadmap")
add_subtitle(slide, "Remaining work and future enhancements")

add_content_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.2),
    "Immediate (This Week)", [
        "Polish cost estimation UI and validation",
        "Add missing field tooltips and help text",
        "Finalize report formatting",
    ], color=PRIMARY)

add_content_card(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.2),
    "Short-term (Next 2 Weeks)", [
        "Database persistence for saved mix designs",
        "Print-to-PDF from report preview",
        "Unit testing for cost estimation tab",
    ], color=SECONDARY)

add_content_card(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4.7),
    "Future Enhancements", [
        "Web version (FastAPI + React)",
        "Additional codes (BS 882, EN 206)",
        "Trial mix comparison tool",
        "Aggregate grading chart visualization",
        "Cloud backup and collaboration",
    ], color=ACCENT)


# ── Slide 15: Q&A ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Questions & Discussion"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "[Your Name]  |  [Email]  |  [GitHub/Portfolio]"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
p3.alignment = PP_ALIGN.CENTER


# ── Save ─────────────────────────────────────────────────────────
output_path = "/home/defy/Documents/projects/civilqntify/CivilQntify_Progress_Slides.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
